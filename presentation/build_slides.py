#!/usr/bin/env python3
"""
Build DND 2026 Datathon presentation.
Run: .venv/bin/python presentation/build_slides.py
Output: presentation/DND_2026_importNumpy.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import sys

# --- Paths ---
ROOT = Path(__file__).parent.parent
FIGURES = ROOT / "reports" / "figures"
OUTPUT = ROOT / "presentation" / "DND_2026_importNumpy.pptx"

# --- Design System ---
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# Colors
NAVY       = RGBColor(0x1B, 0x3A, 0x6B)
CORAL      = RGBColor(0xE8, 0x5D, 0x40)
GREEN      = RGBColor(0x2D, 0x8C, 0x4E)
BLUE_LABEL = RGBColor(0x1B, 0x6B, 0xB0)
PURPLE     = RGBColor(0x7B, 0x5E, 0xA7)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF8)
MID_GRAY   = RGBColor(0x9E, 0x9E, 0x9E)
DARK_GRAY  = RGBColor(0x2D, 0x2D, 0x2D)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# Typography
FONT_FAMILY    = "Calibri"
SIZE_TITLE     = Pt(30)
SIZE_SUBTITLE  = Pt(18)
SIZE_BODY      = Pt(15)
SIZE_CAPTION   = Pt(11)
SIZE_BADGE     = Pt(10)
SIZE_FOOTER    = Pt(9)

# Layout constants
MARGIN_L      = Inches(0.45)
MARGIN_R      = Inches(0.45)
TITLE_BAR_H   = Inches(1.15)
CONTENT_TOP   = Inches(1.25)
CONTENT_H     = Inches(5.75)
COL_MID       = Inches(6.67)
COL_W         = Inches(5.95)
COL_GAP       = Inches(0.30)
FOOTER_TOP    = Inches(7.05)
FOOTER_H      = Inches(0.35)

BADGE_MAP = {
    "DESCRIPTIVE": GREEN,
    "ROBUSTNESS": BLUE_LABEL,
    "ILLUSTRATIVE": PURPLE,
}

# ─────────────────────────── Helpers ───────────────────────────

def _set_font(run, size, bold=False, color=DARK_GRAY, font_name=FONT_FAMILY, italic=False):
    """Apply consistent font settings to a run."""
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    run.font.italic = italic


def add_rect(slide, left, top, width, height, fill_rgb, line=False):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if not line:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a text box with consistent formatting."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    txbox.word_wrap = wrap
    tf = txbox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, font_size, bold=bold, color=color, italic=italic)
    return txbox


def add_image(slide, img_path, left, top, width, height):
    """Add an image; if path missing, add gray placeholder with filename."""
    p = FIGURES / img_path if not Path(img_path).is_absolute() else Path(img_path)
    if p.exists():
        slide.shapes.add_picture(str(p), left, top, width, height)
    else:
        print(f"  WARNING: Missing image {p} — adding placeholder")
        add_rect(slide, left, top, width, height, MID_GRAY)
        add_text_box(slide, f"[Missing: {img_path}]", left, top, width, height,
                     SIZE_CAPTION, color=WHITE, align=PP_ALIGN.CENTER)


def add_evidence_badge(slide, label, left, top):
    """Add a small colored pill badge."""
    color = BADGE_MAP.get(label, MID_GRAY)
    w, h = Inches(1.65), Inches(0.30)
    shape = add_rect(slide, left, top, w, h, color)
    # Add text to the shape
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    _set_font(run, SIZE_BADGE, bold=True, color=WHITE)
    return shape


def add_footer(slide, slide_num, total=12):
    """Add bottom footer: team name left, slide number right."""
    # Team name left
    add_text_box(slide, "importNumpy  |  DND Spring 2026 Datathon",
                 MARGIN_L, FOOTER_TOP, Inches(4), FOOTER_H,
                 SIZE_FOOTER, color=MID_GRAY)
    # Slide number right
    add_text_box(slide, f"{slide_num} of {total}",
                 Inches(11.5), FOOTER_TOP, Inches(1.5), FOOTER_H,
                 SIZE_FOOTER, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def add_title_bar(slide, title_text, bar_color=NAVY, text_color=WHITE):
    """Full-width colored bar across top with white title text."""
    add_rect(slide, 0, 0, SLIDE_W, TITLE_BAR_H, bar_color)
    add_text_box(slide, title_text,
                 MARGIN_L, Inches(0.20), Inches(12.4), Inches(0.80),
                 SIZE_TITLE, bold=True, color=text_color)


def add_trap_header(slide, trap_num):
    """Coral bar with 'TRAP #N' label above title bar."""
    trap_h = Inches(0.42)
    add_rect(slide, 0, 0, SLIDE_W, trap_h, CORAL)
    add_text_box(slide, f"TRAP #{trap_num}",
                 MARGIN_L, Inches(0.05), Inches(3), Inches(0.35),
                 Pt(14), bold=True, color=WHITE)


def add_trap_title_bar(slide, trap_num, title_text):
    """Add trap header + title bar shifted down."""
    add_trap_header(slide, trap_num)
    trap_h = Inches(0.42)
    # Title bar below trap header
    add_rect(slide, 0, trap_h, SLIDE_W, TITLE_BAR_H, NAVY)
    add_text_box(slide, title_text,
                 MARGIN_L, trap_h + Inches(0.20), Inches(12.4), Inches(0.80),
                 SIZE_TITLE, bold=True, color=WHITE)


def add_column_header(slide, text, left, top, width, fill_color):
    """Small colored bar header for a column."""
    h = Inches(0.38)
    add_rect(slide, left, top, width, h, fill_color)
    add_text_box(slide, text, left + Inches(0.1), top + Inches(0.02),
                 width - Inches(0.2), h,
                 SIZE_BADGE, bold=True, color=WHITE)


def add_bullets(slide, lines, left, top, width, height,
                font_size=SIZE_BODY, color=DARK_GRAY, italic=False):
    """Add multi-line bulleted text box."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    txbox.word_wrap = True
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        _set_font(run, font_size, color=color, italic=italic)
    return txbox


def _badge_pos():
    """Default badge position (bottom right)."""
    return Inches(11.23), Inches(6.55)


def add_badges(slide, labels):
    """Add one or more evidence badges stacked at default position."""
    bx, by = _badge_pos()
    for i, label in enumerate(labels):
        add_evidence_badge(slide, label, bx, by + Inches(i * 0.35))


def add_table(slide, headers, rows, left, top, width, height):
    """Add a table with navy header row."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths evenly
    col_w = int(width / n_cols)
    for ci in range(n_cols):
        table.columns[ci].width = col_w

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr
        # Fill header
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '1B3A6B'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                _set_font(run, SIZE_CAPTION, bold=True, color=WHITE)

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    _set_font(run, SIZE_CAPTION, color=DARK_GRAY)

    return table_shape


# ─────────────────────────── Slide Builders ───────────────────────────

def _add_bg(slide):
    """Set light gray background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY


def build_slide_01(prs):
    """Slide 1 — Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_bg(slide)

    # Navy header block
    add_rect(slide, 0, 0, SLIDE_W, Inches(3.8), NAVY)

    # Title
    add_text_box(slide, "5 Conclusions That Would Kill Your Policy",
                 MARGIN_L, Inches(0.6), Inches(12.4), Inches(1.6),
                 Pt(38), bold=True, color=WHITE)

    # Subtitle
    add_text_box(slide, "A Data Integrity Analysis of Global Health Statistics",
                 MARGIN_L, Inches(2.1), Inches(12.4), Inches(0.6),
                 SIZE_SUBTITLE, color=WHITE)

    # Team
    add_text_box(slide, "Team importNumpy",
                 MARGIN_L, Inches(4.2), Inches(5), Inches(0.5),
                 SIZE_BODY, bold=True, color=NAVY)

    # Event
    add_text_box(slide, "DND Spring 2026 Datathon  \u00b7  Cal Poly Pomona",
                 MARGIN_L, Inches(4.7), Inches(8), Inches(0.4),
                 SIZE_CAPTION, color=MID_GRAY)

    # Data note
    add_text_box(slide, "Dataset: Global Health Statistics  \u00b7  1,000,000 rows  \u00b7  20 countries  \u00b7  2000\u20132024",
                 MARGIN_L, Inches(5.2), Inches(10), Inches(0.4),
                 SIZE_CAPTION, color=MID_GRAY)

    add_footer(slide, 1)


def build_slide_02(prs):
    """Slide 2 — What's in the Data"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_title_bar(slide, "What\u2019s in the Data \u2014 and What\u2019s Not")

    # Left column
    col_top = CONTENT_TOP
    add_column_header(slide, "MEASURED \u2713", MARGIN_L, col_top, COL_W, GREEN)

    measured = [
        "\u2022 Outcomes: mortality_rate, recovery_rate, DALYs",
        "\u2022 Access proxies: healthcare_access, doctors_per_1000,",
        "  hospital_beds_per_1000, resource_index",
        "\u2022 Socioeconomic: per_capita_income, education_index",
        "\u2022 Geography proxy: urbanization_rate \u2192 3 tier schemes",
        "\u2022 Structure: country, year, disease, age_group, gender",
        "\u2022 Scale: 1M rows \u00b7 119,978 unique combinations",
    ]
    add_bullets(slide, measured,
                MARGIN_L, col_top + Inches(0.5), COL_W, Inches(4.2))

    # Right column
    right_left = COL_MID + COL_GAP
    add_column_header(slide, "NOT MEASURED \u2717", right_left, col_top, COL_W, CORAL)

    not_measured = [
        "\u2022 Distance-to-care (travel time, proximity)",
        "\u2022 Facility locations or density maps",
        "\u2022 Utilization (visits, admissions, claims)",
        "\u2022 Continuity-of-care metrics",
        "\u2022 Real-world reporting variation or missingness",
    ]
    add_bullets(slide, not_measured,
                right_left, col_top + Inches(0.5), COL_W, Inches(4.2),
                color=CORAL)

    # Bottom note
    add_text_box(slide,
                 "All claims are labeled: DESCRIPTIVE | ROBUSTNESS | ILLUSTRATIVE \u2014 see appendix for evidence contract.",
                 MARGIN_L, Inches(6.3), Inches(10), Inches(0.4),
                 SIZE_CAPTION, color=NAVY, italic=True)

    add_badges(slide, ["DESCRIPTIVE"])
    add_footer(slide, 2)


def build_slide_03(prs):
    """Slide 3 — Clean Data != Realistic Data"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_title_bar(slide, "Clean Data \u2260 Realistic Data \u2014 Six Synthetic Signatures")

    col_top = CONTENT_TOP
    # Left: Data Quality
    add_column_header(slide, "DATA QUALITY \u2713", MARGIN_L, col_top, COL_W, GREEN)
    quality = [
        "\u2022 0% missingness across all 25 columns",
        "\u2022 0 out-of-range values on all rate columns",
        "\u2022 1,000,000 rows, all pass quality flags",
        "\u2022 119,978 / 120,000 possible combinations",
        "  present (99.98%)",
        "\u2022 All 20 disease names correctly mapped",
    ]
    add_bullets(slide, quality,
                MARGIN_L, col_top + Inches(0.5), COL_W, Inches(4.5))

    # Right: Synthetic Signatures
    right_left = COL_MID + COL_GAP
    add_column_header(slide, "SYNTHETIC SIGNATURES \u26a0", right_left, col_top, COL_W, CORAL)
    synth = [
        "\u26a0  Near-Cartesian coverage: 99.98% of all",
        "   (country \u00d7 year \u00d7 disease \u00d7 age \u00d7 gender) cells filled",
        "\u26a0  Near-zero dependence: |r|max = 0.0023",
        "   across 14 variables",
        "\u26a0  Bounded-range match: observed mean/SD match",
        "   Uniform(a,b) expectations to <0.01% error",
        "\u26a0  Homogeneous countries: mortality SD across",
        "   20 countries = 0.013 (ANOVA \u03b7\u00b2 = 0.0000208)",
        "\u26a0  Randomized disease labels: entropy = 3.459",
        "   \u2248 theoretical max",
        "\u26a0  Strict 2-decimal quantization: 100% of",
        "   numeric values",
    ]
    add_bullets(slide, synth,
                right_left, col_top + Inches(0.5), COL_W, Inches(4.5))

    add_badges(slide, ["DESCRIPTIVE"])
    add_footer(slide, 3)


def build_slide_04(prs):
    """Slide 4 — Synthetic Signatures: Headline Metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_title_bar(slide, "Geography Explains 0.002% of Mortality Variance")

    # Left: Table
    headers = ["Metric", "This Dataset", "Real-World Expected"]
    rows = [
        ["Access \u2194 Mortality (r)", "0.00008", "\u22120.30 to \u22120.50"],
        ["Income \u2194 Mortality (r)", "\u22120.0015", "\u22120.20 to \u22120.40"],
        ["Country ANOVA \u03b7\u00b2", "0.0000208", "0.05 to 0.30"],
        ["Year trend \u03b7\u00b2", "0.0000106", "0.01 to 0.15"],
    ]
    add_table(slide, headers, rows,
              MARGIN_L, CONTENT_TOP, COL_W, Inches(2.5))

    # Caption below table
    add_text_box(slide,
                 "Expected ranges from WHO/Lancet global health literature. Observed values are descriptive of this dataset only.",
                 MARGIN_L, CONTENT_TOP + Inches(2.7), COL_W, Inches(0.6),
                 SIZE_CAPTION, color=MID_GRAY, italic=True)

    # Right: Image
    add_image(slide, "bench_q6_dumbbell_correlations.png",
              COL_MID + COL_GAP, CONTENT_TOP, COL_W, Inches(5.2))

    add_badges(slide, ["DESCRIPTIVE", "ILLUSTRATIVE"])
    add_footer(slide, 4)


def build_slide_05(prs):
    """Slide 5 — Trap #1: Rural Outcomes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_trap_title_bar(slide, 1, '\u201cRural Areas Have Worse Outcomes\u201d \u2014 Evidence: ABSENT')

    trap_offset = Inches(0.42)
    ct = CONTENT_TOP + trap_offset

    lines = [
        "Common claim:",
        "  Urban/rural classification predicts health outcomes.",
        "",
        "What we found (DESCRIPTIVE):",
        "  \u2022 Cohen\u2019s d < 0.01 for all tier pairs (negligible effect)",
        "  \u2022 ANOVA \u03b7\u00b2 = 0.0000208 (geography explains 0.002%)",
        "  \u2022 p < 0.001 \u2014 but N=1M makes everything significant",
        "  \u2022 Absolute mortality difference Rural vs Urban: < 0.05pp",
        "",
        "Why it\u2019s premature:",
        "  Statistical significance \u2260 practical significance.",
        "  With 1M rows, any noise clears p-value thresholds.",
        "",
        "Safer framing:",
        "  \u201cWe cannot distinguish outcomes across urbanization",
        "  tiers with this dataset. Report effect sizes, not p-values.\u201d",
    ]
    add_bullets(slide, lines, MARGIN_L, ct, COL_W, Inches(4.8))

    add_image(slide, "q1_effect_size_dotplot.png",
              COL_MID + COL_GAP, ct, COL_W, Inches(5.0))

    # Green note
    add_text_box(slide, "\u2713 Use effect sizes (Cohen\u2019s d, \u03b7\u00b2), not p-values with large N",
                 MARGIN_L, Inches(6.4), COL_W, Inches(0.35),
                 SIZE_CAPTION, color=GREEN)

    add_badges(slide, ["DESCRIPTIVE"])
    add_footer(slide, 5)


def build_slide_06(prs):
    """Slide 6 — Trap #2: Access → Outcomes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_trap_title_bar(slide, 2, '\u201cHealthcare Access Improves Outcomes\u201d \u2014 Evidence: FLAT')

    trap_offset = Inches(0.42)
    ct = CONTENT_TOP + trap_offset

    lines = [
        "Common claim:",
        "  Higher healthcare access scores \u2192 lower mortality rates.",
        "",
        "What we found (DESCRIPTIVE):",
        "  \u2022 Pearson r (access \u2194 mortality) = 0.00008",
        "  \u2022 R\u00b2 < 0.000001 \u2014 access explains essentially 0% of variance",
        "  \u2022 No dose-response at any access decile threshold",
        "  \u2022 Spearman \u03c1 consistent: no monotonic relationship",
        "",
        "Critical gap:",
        "  Distance-to-care and utilization are NOT in this dataset.",
        "  \u201cHealthcare access\u201d is a proxy index, not a distance measure.",
        "",
        "Safer framing:",
        "  \u201cThis dataset cannot support an access-outcome causal claim.",
        "  Collecting travel time and utilization data is the prerequisite.\u201d",
    ]
    add_bullets(slide, lines, MARGIN_L, ct, COL_W, Inches(4.8))

    add_image(slide, "q2_scatter_access_mortality.png",
              COL_MID + COL_GAP, ct, COL_W, Inches(5.0))

    # Warning note
    add_text_box(slide,
                 "\u26a0 Distance-to-care: NOT MEASURED. Access proxies used only. See prompt_gap_matrix.md.",
                 MARGIN_L, Inches(6.4), Inches(10), Inches(0.35),
                 SIZE_CAPTION, color=CORAL)

    add_badges(slide, ["DESCRIPTIVE"])
    add_footer(slide, 6)


def build_slide_07(prs):
    """Slide 7 — Trap #3: Country Outliers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_trap_title_bar(slide, 3, '\u201cCountry X Is an Outlier Worth Studying\u201d \u2014 Evidence: NOISE')

    trap_offset = Inches(0.42)
    ct = CONTENT_TOP + trap_offset

    lines = [
        "Common claim:",
        "  High z-score countries represent exceptional systems",
        "  worth studying or emulating.",
        "",
        "What we found (DESCRIPTIVE):",
        "  \u2022 Some countries show high z-scores on access vs outcomes",
        "  \u2022 Absolute mortality differences: < 0.5 percentage points",
        "  \u2022 Between-country SD: 0.013 on mortality (scale: 0\u201310)",
        "  \u2022 Effect is noise-level at this resolution",
        "",
        "Why it\u2019s premature:",
        "  Z-scores inflate when within-group variance is near-zero.",
        "  A 0.5pp absolute difference is not a policy signal.",
        "",
        "Safer framing:",
        "  \u201cNo country deviates meaningfully in absolute terms.",
        "  Statistical outliers here are artifacts of near-zero variance,",
        "  not evidence of a distinctive health system.\u201d",
    ]
    add_bullets(slide, lines, MARGIN_L, ct, COL_W, Inches(4.8))

    add_image(slide, "q3_quadrant_scatter.png",
              COL_MID + COL_GAP, ct, COL_W, Inches(5.0))

    add_badges(slide, ["DESCRIPTIVE"])
    add_footer(slide, 7)


def build_slide_08(prs):
    """Slide 8 — Trap #4: Tier Definitions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_trap_title_bar(slide, 4, '\u201cTier Definitions Change Conclusions\u201d \u2014 Evidence: IRRELEVANT')

    trap_offset = Inches(0.42)
    ct = CONTENT_TOP + trap_offset

    lines = [
        "Common claim:",
        "  How you define \u201curban\u201d vs \u201crural\u201d changes your findings.",
        "",
        "What we found (ROBUSTNESS):",
        "  \u2022 3 tier schemes tested: Default 3-tier, Binary, 4-tier",
        "  \u2022 Maximum spread across schemes: < 0.04 percentage points",
        "  \u2022 ANOVA \u03b7\u00b2 \u2248 0 under all three definitions",
        "  \u2022 Conclusions are identical regardless of classification choice",
        "",
        "Why this matters:",
        "  Robustness is confirmed \u2014 but only because there is",
        "  no signal to be sensitive to in the first place.",
        "  Tier choice is genuinely irrelevant when the underlying",
        "  effect size is zero.",
        "",
        "Safer framing:",
        "  \u201cOur robustness check is valid, but it validates",
        "  a null result \u2014 not a policy-relevant pattern.\u201d",
    ]
    add_bullets(slide, lines, MARGIN_L, ct, COL_W, Inches(4.8))

    add_image(slide, "q4_3panel_tier_schemes.png",
              COL_MID + COL_GAP, ct, COL_W, Inches(5.0))

    add_badges(slide, ["ROBUSTNESS"])
    add_footer(slide, 8)


def build_slide_09(prs):
    """Slide 9 — Trap #5: Sparse Reporting"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_trap_title_bar(slide, 5, '\u201cSparse Reporting Drives Uncertainty Here\u201d \u2014 Evidence: NOT OBSERVED')

    trap_offset = Inches(0.42)
    ct = CONTENT_TOP + trap_offset

    lines = [
        "Common claim:",
        "  Data-sparse regions show high uncertainty that distorts",
        "  conclusions about access and outcomes.",
        "",
        "What we found (DESCRIPTIVE):",
        "  \u2022 Missingness: 0% across all 25 columns",
        "  \u2022 All 640 subgroups (country \u00d7 disease \u00d7 age \u00d7 gender):",
        "    count \u2265 30 \u2014 all above minimum threshold",
        "  \u2022 Confidence intervals are uniformly narrow",
        "",
        "Simulation (ILLUSTRATIVE):",
        "  We demonstrate what sparsity would look like via",
        "  subsampling: removing 95% of data widens CIs dramatically.",
        "  This is what the data WOULD show if sparse \u2014 it doesn\u2019t.",
        "",
        "Safer framing:",
        "  \u201cSparse reporting is a real-world concern we cannot",
        "  assess from this dataset. The dataset is fully populated",
        "  by design. We model the risk via simulation only.\u201d",
    ]
    add_bullets(slide, lines, MARGIN_L, ct, COL_W, Inches(4.8))

    add_image(slide, "q5_forest_full_vs_sparse.png",
              COL_MID + COL_GAP, ct, COL_W, Inches(5.0))

    add_badges(slide, ["DESCRIPTIVE", "ILLUSTRATIVE"])
    add_footer(slide, 9)


def build_slide_10(prs):
    """Slide 10 — Robustness: Repeated-Cell Check"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_title_bar(slide, "Results Hold Under Repeated-Cell Robustness Check")

    # Intro text
    add_text_box(slide,
                 "Each (country \u00d7 year \u00d7 disease \u00d7 age \u00d7 gender) cell appears ~8.3 times in the dataset. "
                 "We compared base views vs. dedup-at-cell-grain views to confirm repeated sampling does not inflate any finding.",
                 MARGIN_L, CONTENT_TOP, Inches(12.4), Inches(0.7),
                 SIZE_BODY, color=DARK_GRAY)

    # Table
    headers = ["Question", "Base Value", "Deduped Value", "Delta", "Conclusion"]
    rows = [
        ["Q1 Mortality tier gap", "0.0492 pp", "0.0489 pp", "0.0003 pp", "No artifact"],
        ["Q2 Access\u2194mortality r", "0.00008", "0.00007", "0.00001", "No artifact"],
        ["Q3 Max country z-score", "2.31", "2.28", "0.03", "No artifact"],
        ["Q4 Tier scheme spread", "0.038 pp", "0.037 pp", "0.001 pp", "No artifact"],
        ["Q5 Min group count", "8,318", "1", "\u2014", "By design"],
        ["Q6 Max |r| any pair", "0.0023", "0.0021", "0.0002", "No artifact"],
    ]
    add_table(slide, headers, rows,
              MARGIN_L, Inches(2.2), Inches(12.4), Inches(3.5))

    # Bottom note
    add_text_box(slide,
                 "\u2713 All deltas are noise-level. Repeated sampling per cell does not create spurious patterns.",
                 MARGIN_L, Inches(6.0), Inches(10), Inches(0.4),
                 SIZE_CAPTION, color=GREEN)

    add_badges(slide, ["ROBUSTNESS"])
    add_footer(slide, 10)


def build_slide_11(prs):
    """Slide 11 — Expected vs Observed"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_title_bar(slide, "This Is Why Policy Prescriptions Would Be Premature")

    # Left column
    col_top = Inches(1.4)
    add_column_header(slide, "REAL-WORLD BENCHMARK (Illustrative)", MARGIN_L, col_top, COL_W, NAVY)
    add_image(slide, "bench_q2_scatter_comparison.png",
              MARGIN_L, Inches(1.9), COL_W, Inches(4.3))
    add_text_box(slide, "Literature-based: access \u2194 mortality r \u2248 \u22120.45",
                 MARGIN_L, Inches(6.3), COL_W, Inches(0.3),
                 SIZE_CAPTION, color=MID_GRAY)

    # Right column
    right_left = COL_MID + COL_GAP
    add_column_header(slide, "THIS DATASET (Descriptive)", right_left, col_top, COL_W, CORAL)
    add_image(slide, "q2_scatter_access_mortality.png",
              right_left, Inches(1.9), COL_W, Inches(4.3))
    add_text_box(slide, "Observed: r = 0.00008 \u2014 effectively zero",
                 right_left, Inches(6.3), COL_W, Inches(0.3),
                 SIZE_CAPTION, color=CORAL)

    # Bottom center
    add_text_box(slide,
                 "The gap between benchmark and observed is the reason policy prescriptions from this dataset would be premature.",
                 Inches(1.5), Inches(6.65), Inches(10.3), Inches(0.5),
                 SIZE_BODY, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_badges(slide, ["ILLUSTRATIVE", "DESCRIPTIVE"])
    add_footer(slide, 11)


def build_slide_12(prs):
    """Slide 12 — Guardrails + Action Plan"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide)
    add_title_bar(slide, "What Responsible Analysis Requires Before Policy Action")

    col_top = Inches(1.4)
    # Left: Data Readiness Scorecard
    add_column_header(slide, "DATA READINESS SCORECARD", MARGIN_L, col_top, COL_W, NAVY)
    add_image(slide, "bench_data_readiness_scorecard.png",
              MARGIN_L, Inches(1.9), COL_W, Inches(4.0))
    add_text_box(slide, "Traffic-light assessment: 8 data requirements for responsible policy analysis",
                 MARGIN_L, Inches(6.0), COL_W, Inches(0.3),
                 SIZE_CAPTION, color=MID_GRAY)

    # Right: Roadmap
    right_left = COL_MID + COL_GAP
    add_column_header(slide, "30 / 60 / 90-DAY ROADMAP", right_left, col_top, COL_W, GREEN)
    add_image(slide, "bench_action_timeline.png",
              right_left, Inches(1.9), COL_W, Inches(4.0))
    add_text_box(slide, "12 concrete next steps \u2014 deployable on real data when collected",
                 right_left, Inches(6.0), COL_W, Inches(0.3),
                 SIZE_CAPTION, color=MID_GRAY)

    # Bottom quote
    add_text_box(slide,
                 '\u201cThe most valuable analysis is the one that knows when NOT to make a recommendation \u2014 '
                 'and shows exactly what it WOULD recommend with the right data.\u201d',
                 Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.5),
                 SIZE_CAPTION, italic=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_badges(slide, ["ILLUSTRATIVE"])
    add_footer(slide, 12)


# ─────────────────────────── Main ───────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.50)

    builders = [
        build_slide_01, build_slide_02, build_slide_03, build_slide_04,
        build_slide_05, build_slide_06, build_slide_07, build_slide_08,
        build_slide_09, build_slide_10, build_slide_11, build_slide_12,
    ]

    for i, builder in enumerate(builders, 1):
        print(f"  Building slide {i}/12...")
        builder(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"\nSaved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
